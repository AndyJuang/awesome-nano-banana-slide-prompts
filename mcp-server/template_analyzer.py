"""
template_analyzer.py
分析使用者上傳的 PPTX 模版，提取風格規格供後續生成使用。

提取內容：
  - 投影片尺寸與比例
  - 色彩調色盤（依出現頻率排序）
  - 字型名稱與字號分布
  - 各頁版型推測（封面/條列/表格/圖文等）
  - 佈局關鍵尺寸（標題位置、邊距、色塊比例）
"""
from pptx import Presentation
from pptx.util import Inches
from collections import Counter
import json


# ── 內部工具 ────────────────────────────────────────────────────────────────

def _rgb_to_hex(rgb) -> str | None:
    try:
        return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
    except Exception:
        return None


def _shape_fill_hex(shape) -> str | None:
    try:
        if shape.fill.type == 1:            # MSO_FILL.SOLID
            return _rgb_to_hex(shape.fill.fore_color.rgb)
    except Exception:
        pass
    return None


def _classify_slide(shapes: list, slide_idx: int, total: int) -> str:
    """根據形狀數量與位置推測版型類型"""
    n = len(shapes)
    fills = [s.get("fill") for s in shapes if s.get("fill")]

    # 第一頁 → 封面
    if slide_idx == 0:
        return "cover"
    # 最後一頁 → 結尾
    if slide_idx == total - 1:
        return "closing"
    # 形狀少 → 章節分隔
    if n <= 6:
        return "section_divider"
    # 有大型寬矩形（可能是表格頁）
    wide_shapes = [s for s in shapes
                   if s["pos"]["width"] > 10 and s["pos"]["height"] < 1.5]
    if wide_shapes:
        return "table"
    # 有圖片位置（佔左側或右側 40%+）
    left_heavy = any(
        s["pos"]["left"] < 1 and s["pos"]["width"] > 5
        for s in shapes
    )
    if left_heavy and n < 12:
        return "image_split"
    # 預設為條列式
    return "bullet_points"


# ── 主要分析函式 ─────────────────────────────────────────────────────────────

def analyze_pptx_template(pptx_path: str) -> dict:
    """
    分析 PPTX 模版，回傳標準化風格規格 dict。

    回傳欄位：
      dimensions        投影片尺寸
      color_palette     色彩調色盤（hex list，依頻率排序）
      primary_color     主色
      accent_color      強調色（第二常用）
      bg_color          最常見背景色
      fonts             字型名稱清單
      heading_size      推測標題字號
      body_size         推測內文字號
      slide_count       總頁數
      layout_map        每頁推測版型
      key_elements      關鍵設計元素描述
      style_prompt      給 AI 生成時使用的風格描述字串
    """
    prs = Presentation(pptx_path)

    color_counter: Counter = Counter()
    font_counter:  Counter = Counter()
    size_counter:  Counter = Counter()
    bg_colors:     list    = []
    slide_infos:   list    = []

    total = len(prs.slides)

    for idx, slide in enumerate(prs.slides):
        # 背景色
        try:
            bg = slide.background.fill
            if bg.type == 1:
                hex_c = _rgb_to_hex(bg.fore_color.rgb)
                if hex_c:
                    bg_colors.append(hex_c)
                    color_counter[hex_c] += 5    # 背景色給予較高權重
        except Exception:
            pass

        shapes_data = []
        for shape in slide.shapes:
            pos = {
                "left":   round(shape.left   / 914400, 2),
                "top":    round(shape.top    / 914400, 2),
                "width":  round(shape.width  / 914400, 2),
                "height": round(shape.height / 914400, 2),
            }
            s = {"pos": pos}

            # 填色
            hex_fill = _shape_fill_hex(shape)
            if hex_fill:
                s["fill"] = hex_fill
                color_counter[hex_fill] += 1

            # 文字
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            font_counter[run.font.name] += 1
                        if run.font.size:
                            size_counter[round(run.font.size.pt)] += 1
                        try:
                            if run.font.color.type:
                                hex_c = _rgb_to_hex(run.font.color.rgb)
                                if hex_c:
                                    color_counter[hex_c] += 1
                        except Exception:
                            pass

            shapes_data.append(s)

        slide_infos.append({
            "slide": idx + 1,
            "shapes": shapes_data,
            "inferred_layout": _classify_slide(shapes_data, idx, total),
        })

    # ── 彙整 ──────────────────────────────────────────────────────────────
    palette = [c for c, _ in color_counter.most_common(12)]
    # 濾掉純白純黑（通常不是品牌色）
    brand_colors = [c for c in palette
                    if c not in ("#FFFFFF", "#000000", "#ffffff", "#000000")]

    primary_color = brand_colors[0] if brand_colors else "#1B4F9B"
    accent_color  = brand_colors[1] if len(brand_colors) > 1 else "#0091D5"
    bg_color      = Counter(bg_colors).most_common(1)[0][0] if bg_colors else "#FFFFFF"

    fonts = [f for f, _ in font_counter.most_common(5)]
    sizes = dict(size_counter.most_common(10))

    # 推測標題 / 內文字號
    all_sizes = sorted(sizes.keys(), reverse=True)
    heading_size = all_sizes[0] if all_sizes else 24
    body_size    = all_sizes[2] if len(all_sizes) > 2 else 14

    layout_map = [
        {"slide": s["slide"], "layout": s["inferred_layout"]}
        for s in slide_infos
    ]

    # 關鍵設計元素
    key_elements = []
    first = slide_infos[0]["shapes"] if slide_infos else []
    if any(s["pos"]["width"] > 6 and s.get("fill") for s in first):
        key_elements.append("large color block (brand panel)")
    if any(s["pos"]["width"] < 4 and s["pos"]["top"] > 5 for s in first):
        key_elements.append("bottom accent bar")
    if len(first) > 8:
        key_elements.append("photo grid / multi-image layout")

    # 風格描述字串（供 Gemini prompt 使用）
    style_prompt = (
        f"Corporate presentation style. "
        f"Primary brand color: {primary_color}, accent: {accent_color}, "
        f"background: {bg_color}. "
        f"Typography: {', '.join(fonts[:2]) if fonts else 'sans-serif'}, "
        f"heading {heading_size}pt / body {body_size}pt. "
        f"Design elements: {', '.join(key_elements) if key_elements else 'clean minimal'}."
    )

    return {
        "dimensions": {
            "width_inches":  round(prs.slide_width.inches,  2),
            "height_inches": round(prs.slide_height.inches, 2),
        },
        "color_palette":  palette,
        "primary_color":  primary_color,
        "accent_color":   accent_color,
        "bg_color":       bg_color,
        "fonts":          fonts,
        "heading_size":   heading_size,
        "body_size":      body_size,
        "slide_count":    total,
        "layout_map":     layout_map,
        "key_elements":   key_elements,
        "style_prompt":   style_prompt,
    }
