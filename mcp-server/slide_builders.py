"""
slide_builders.py
各版型的 python-pptx 實作函式
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

SLIDE_WIDTH  = Inches(13.33)
SLIDE_HEIGHT = Inches(7.50)

CHAPTER_LABELS = [
    "ONE","TWO","THREE","FOUR","FIVE",
    "SIX","SEVEN","EIGHT","NINE","TEN",
]

# ── 色彩工具 ────────────────────────────────────────────────────────────────

def parse_color(val, default=(0x1B, 0x4F, 0x9B)) -> RGBColor:
    """接受 '#RRGGBB'、(r,g,b) tuple 或 RGBColor"""
    if isinstance(val, RGBColor):
        return val
    if isinstance(val, str):
        h = val.lstrip("#")
        if len(h) == 6:
            return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))
    if isinstance(val, (list, tuple)) and len(val) == 3:
        return RGBColor(*val)
    return RGBColor(*default)

# ── 基礎元素 ────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height,
             fill_color=None, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(1, int(left), int(top), int(width), int(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=14, bold=False, color=None,
             align=PP_ALIGN.LEFT, font_name="Microsoft JhengHei"):
    box = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    if color:
        run.font.color.rgb = color
    return box

# ── 版型實作 ─────────────────────────────────────────────────────────────────

def _build_cover_or_closing(slide, s: dict):
    """封面頁 & 結尾頁（共用結構，is_closing 由 layout 欄位判斷）"""
    is_closing   = s.get("layout") == "closing"
    brand_color  = parse_color(s.get("brand_color", "#1B4F9B"))
    logo_text    = s.get("logo_text", "{brand_name}")
    image_paths  = s.get("image_paths") or [None, None, None]

    # — 左側三塊照片佔位色 —
    photo_specs = [
        (0.00, 0.00, 4.20, 4.50, 0x19, 0x32, 0x5A),
        (4.35, 0.00, 2.05, 2.15, 0x0F, 0x3C, 0x6E),
        (4.35, 2.35, 2.05, 2.15, 0x0A, 0x23, 0x4B),
    ]
    for i, (l, t, w, h, r, g, b) in enumerate(photo_specs):
        img = image_paths[i] if i < len(image_paths) else None
        if img and Path(img).exists():
            slide.shapes.add_picture(img, Inches(l), Inches(t), Inches(w), Inches(h))
        else:
            add_rect(slide, Inches(l), Inches(t), Inches(w), Inches(h),
                     fill_color=RGBColor(r, g, b))

    # — 右側品牌色塊 —
    add_rect(slide, Inches(6.50), 0, Inches(6.83), Inches(4.50), fill_color=brand_color)

    # — 左上角 Logo —
    add_text(slide, logo_text, Inches(0.20), Inches(0.10), Inches(6.00), Inches(0.45),
             font_size=9, bold=True, color=RGBColor(255,255,255), font_name="Calibri")

    # — 品牌名（封面才顯示）—
    if not is_closing:
        for line, top in [
            (s.get("brand_name_line_1",""), 0.70),
            (s.get("brand_name_line_2",""), 1.85),
        ]:
            if line:
                add_text(slide, line, Inches(6.60), Inches(top),
                         Inches(6.60), Inches(1.10),
                         font_size=46, bold=True,
                         color=RGBColor(255,255,255), font_name="Calibri")

    # — 下半白底 —
    add_rect(slide, 0, Inches(4.55), SLIDE_WIDTH, Inches(2.95),
             fill_color=RGBColor(255,255,255))

    # — 主標題 —
    main_title = s.get("closing_primary","感謝您") if is_closing else s.get("main_title","")
    main_size  = 50 if is_closing else 38
    add_text(slide, main_title, Inches(0.80), Inches(4.65),
             Inches(11.70), Inches(1.10),
             font_size=main_size, bold=True, color=RGBColor(64,64,64))

    # — 副標題 —
    subtitle  = s.get("closing_secondary","Thank you") if is_closing else s.get("subtitle","")
    sub_size  = 22 if is_closing else 17
    sub_top   = 5.75 if is_closing else 5.85
    if subtitle:
        add_text(slide, subtitle, Inches(0.80), Inches(sub_top),
                 Inches(11.70), Inches(0.60),
                 font_size=sub_size, bold=True,
                 color=RGBColor(64,64,64), font_name="Calibri")

    # — CTA 按鈕 —
    cta   = "END" if is_closing else s.get("cta_text","START")
    btn_x = Inches(5.70) if is_closing else Inches(5.85)
    btn_w = Inches(1.90) if is_closing else Inches(1.60)
    add_rect(slide, btn_x, Inches(6.72), btn_w, Inches(0.45),
             fill_color=RGBColor(255,255,255), border_color=RGBColor(180,180,180))
    add_text(slide, cta, btn_x, Inches(6.74), btn_w, Inches(0.38),
             font_size=12 if is_closing else 11,
             align=PP_ALIGN.CENTER, font_name="Calibri")

    # — 夥伴/備注文字 —
    partner = s.get("partner_text","")
    if partner:
        add_text(slide, partner, Inches(10.50), Inches(6.90),
                 Inches(2.70), Inches(0.38),
                 font_size=9, align=PP_ALIGN.RIGHT, font_name="Calibri")


def _build_contents(slide, s: dict):
    """深色左欄數字徽章目錄"""
    panel_color  = parse_color(s.get("panel_color",     "#2D415F"))
    badge_odd    = parse_color(s.get("badge_color_odd",  "#1B4F9B"))
    badge_even   = parse_color(s.get("badge_color_even", "#808080"))
    logo_text    = s.get("logo_text", "{brand_name}")

    # 左欄
    add_rect(slide, 0, 0, Inches(5.80), SLIDE_HEIGHT, fill_color=panel_color)
    add_rect(slide, Inches(0.35), Inches(2.70), Inches(0.32), Inches(0.32),
             fill_color=RGBColor(0x1B,0x4F,0x9B))
    add_rect(slide, Inches(0.75), Inches(2.90), Inches(0.22), Inches(0.22),
             fill_color=RGBColor(0x00,0x91,0xD5))

    add_text(slide, s.get("label_primary","目錄"),
             Inches(0.50), Inches(2.85), Inches(4.50), Inches(0.85),
             font_size=34, bold=True, color=RGBColor(255,255,255))
    add_text(slide, s.get("label_secondary","contents"),
             Inches(0.50), Inches(3.65), Inches(4.50), Inches(0.60),
             font_size=20, bold=True, color=RGBColor(255,255,255), font_name="Calibri")

    # 目錄項目
    BADGE_W = Inches(0.58)
    for i, item in enumerate(s.get("items",[])[:6]):
        y = Inches(0.35) + Inches(1.12) * i
        bc = badge_odd if i % 2 == 0 else badge_even

        add_rect(slide, Inches(6.10), y, BADGE_W, BADGE_W, fill_color=bc)
        add_text(slide, str(i+1),
                 Inches(6.10), y + Inches(0.08), BADGE_W, Inches(0.45),
                 font_size=18, bold=True, color=RGBColor(255,255,255),
                 align=PP_ALIGN.CENTER, font_name="Calibri")

        add_text(slide, item.get("title",""),
                 Inches(6.10) + BADGE_W + Inches(0.22), y,
                 Inches(6.00), Inches(0.38),
                 font_size=16, bold=True, color=RGBColor(64,64,64))

        if item.get("subtitle"):
            add_text(slide, item["subtitle"],
                     Inches(6.10) + BADGE_W + Inches(0.22), y + Inches(0.37),
                     Inches(6.00), Inches(0.26),
                     font_size=10, color=RGBColor(128,128,128), font_name="Calibri")

        add_rect(slide, Inches(6.10), y + Inches(0.68),
                 Inches(6.80), Inches(0.025),
                 fill_color=RGBColor(240,240,240))

    # 右下 Logo
    add_text(slide, logo_text, Inches(11.30), Inches(7.10),
             Inches(1.90), Inches(0.32),
             font_size=9, bold=True, color=RGBColor(0x1B,0x4F,0x9B), font_name="Calibri")


def _build_section_divider(slide, s: dict):
    """PART 徽章章節分隔頁"""
    badge_color  = parse_color(s.get("brand_color",  "#1B4F9B"))
    accent_color = parse_color(s.get("accent_color", "#0091D5"))
    logo_text    = s.get("logo_text", "{brand_name}")
    part_number  = int(s.get("part_number", 1))
    section_title = s.get("section_title", "")

    # 淺灰背景
    add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT,
             fill_color=RGBColor(0xF8,0xF8,0xF8))

    # 右上 Logo
    add_text(slide, logo_text, Inches(11.30), Inches(0.08),
             Inches(1.90), Inches(0.45),
             font_size=9, bold=True, color=badge_color, font_name="Calibri")

    # PART 徽章
    add_rect(slide, Inches(1.40), Inches(1.70), Inches(2.60), Inches(2.70),
             fill_color=badge_color)
    add_text(slide, "P", Inches(1.50), Inches(1.80), Inches(1.00), Inches(1.30),
             font_size=68, bold=True, color=RGBColor(255,255,255))
    add_text(slide, f"ART {part_number:02d}",
             Inches(2.40), Inches(3.10), Inches(1.50), Inches(0.55),
             font_size=20, bold=True, color=RGBColor(255,255,255), font_name="Calibri")

    # 空心輪廓方塊
    add_rect(slide, Inches(1.65), Inches(4.55), Inches(1.50), Inches(1.50),
             fill_color=RGBColor(255,255,255), border_color=badge_color)

    # 右側標題
    add_rect(slide, Inches(4.70), Inches(2.95), Inches(0.25), Inches(0.25),
             fill_color=accent_color)
    add_text(slide, section_title, Inches(5.10), Inches(2.68),
             Inches(7.50), Inches(0.85),
             font_size=34, bold=True, color=RGBColor(64,64,64))
    add_rect(slide, Inches(5.10), Inches(3.60), Inches(6.80), Inches(0.06),
             fill_color=accent_color)

    # 右下角英文章節標籤
    label = (CHAPTER_LABELS[part_number-1]
             if 1 <= part_number <= len(CHAPTER_LABELS)
             else str(part_number))
    add_rect(slide, Inches(11.95), Inches(6.82), Inches(1.38), Inches(0.50),
             fill_color=badge_color)
    add_text(slide, label, Inches(11.95), Inches(6.84),
             Inches(1.38), Inches(0.42),
             font_size=13, bold=True, color=RGBColor(255,255,255),
             align=PP_ALIGN.CENTER, font_name="Calibri")


def _build_bullet_points(slide, s: dict):
    """階層式條列內容頁（兩層：分組標題 + 子項目）"""
    accent_color = parse_color(s.get("accent_color", "#0091D5"))
    logo_text    = s.get("logo_text", "{brand_name}")
    slide_title  = s.get("slide_title", "")
    groups       = s.get("groups", [])

    NAVY  = RGBColor(0x1B,0x4F,0x9B)
    DARK  = RGBColor(64,64,64)
    WHITE = RGBColor(255,255,255)

    # 右上裝飾
    add_text(slide, logo_text, Inches(11.30), Inches(0.08),
             Inches(1.90), Inches(0.45),
             font_size=9, bold=True, color=NAVY, font_name="Calibri")
    add_rect(slide, Inches(12.25), Inches(0.22), Inches(0.38), Inches(0.38), fill_color=NAVY)
    add_rect(slide, Inches(11.90), Inches(0.45), Inches(0.24), Inches(0.24), fill_color=accent_color)
    add_rect(slide, Inches(12.58), Inches(0.52), Inches(0.19), Inches(0.19),
             fill_color=WHITE, border_color=RGBColor(180,180,180))

    # 標題列
    add_rect(slide, Inches(0.38), Inches(0.75), Inches(0.22), Inches(0.22), fill_color=accent_color)
    add_text(slide, slide_title, Inches(0.72), Inches(0.62),
             Inches(11.00), Inches(0.58),
             font_size=24, bold=True, color=DARK)
    add_rect(slide, Inches(0.38), Inches(1.28), Inches(12.00), Inches(0.04), fill_color=accent_color)

    # 內容區
    y = Inches(1.50)
    for group in groups[:3]:
        add_rect(slide, Inches(0.38), y, Inches(0.16), Inches(0.30), fill_color=accent_color)
        add_text(slide, group.get("title",""), Inches(0.66), y,
                 Inches(12.00), Inches(0.42),
                 font_size=15, bold=True, color=DARK)
        y += Inches(0.57)

        for bullet in group.get("bullets",[])[:5]:
            add_rect(slide, Inches(0.55), y + Inches(0.15),
                     Inches(0.10), Inches(0.10), fill_color=accent_color)
            add_text(slide, bullet, Inches(0.78), y,
                     Inches(12.00), Inches(0.40),
                     font_size=14, color=DARK)
            y += Inches(0.44)
        y += Inches(0.50)

    # 底部裝飾線
    add_rect(slide, Inches(0.38), Inches(7.18), Inches(12.60), Inches(0.05),
             fill_color=RGBColor(0xE8,0xF4,0xFD))


def _build_table(slide, s: dict):
    """基礎表格頁"""
    accent_color = parse_color(s.get("accent_color", "#0091D5"))
    light_row    = parse_color(s.get("even_row_color", "#E8F4FD"))
    logo_text    = s.get("logo_text", "{brand_name}")
    slide_title  = s.get("slide_title", "")
    col_headers  = s.get("col_headers", [])
    rows_data    = s.get("rows_data", [])

    NAVY  = RGBColor(0x1B,0x4F,0x9B)
    DARK  = RGBColor(64,64,64)
    WHITE = RGBColor(255,255,255)

    add_text(slide, logo_text, Inches(11.30), Inches(0.08),
             Inches(1.90), Inches(0.45),
             font_size=9, bold=True, color=NAVY, font_name="Calibri")
    add_rect(slide, Inches(0.38), Inches(0.45), Inches(0.12), Inches(0.50),
             fill_color=accent_color)
    add_text(slide, slide_title, Inches(0.58), Inches(0.40),
             SLIDE_WIDTH - Inches(1.00), Inches(0.60),
             font_size=22, bold=True, color=DARK)

    if not col_headers or not rows_data:
        return

    tbl = slide.shapes.add_table(
        len(rows_data) + 1, len(col_headers),
        Inches(0.40), Inches(1.30),
        SLIDE_WIDTH - Inches(0.80), SLIDE_HEIGHT - Inches(1.80),
    ).table

    for ci, header in enumerate(col_headers):
        cell = tbl.cell(0, ci)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = accent_color
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.bold  = True
        run.font.color.rgb = WHITE
        run.font.size  = Pt(14)

    for ri, row in enumerate(rows_data):
        bg = light_row if ri % 2 == 1 else WHITE
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size  = Pt(13)
            run.font.color.rgb = DARK


# ── 進階版型建構函式 ──────────────────────────────────────────────────────────

def _build_timeline(slide, s: dict):
    """橫向時間軸 / 里程碑"""
    accent  = parse_color(s.get("accent_color", "#0091D5"))
    primary = parse_color(s.get("brand_color",  "#1B4F9B"))
    DARK    = RGBColor(64,64,64)
    WHITE   = RGBColor(255,255,255)
    title   = s.get("slide_title","")
    nodes   = s.get("nodes", [])   # list of {date, title, description}

    # 頁面標題
    add_text(slide, title, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             font_size=24, bold=True, color=DARK)

    # 橫向時間軸主線
    line_y = Inches(3.8)
    add_rect(slide, Inches(0.8), line_y, Inches(11.7), Inches(0.06),
             fill_color=accent)

    n = len(nodes[:6])
    if n == 0:
        return
    step = Inches(11.7) / max(n - 1, 1)

    for i, node in enumerate(nodes[:6]):
        x = Inches(0.8) + step * i
        is_above = (i % 2 == 0)   # 奇數節點標籤在上，偶數在下

        # 節點圓點（用正方形代替）
        dot_size = Inches(0.22)
        add_rect(slide, x - dot_size/2, line_y - dot_size/2,
                 dot_size, dot_size, fill_color=primary)

        # 垂直連接線
        connector_h = Inches(1.0)
        if is_above:
            add_rect(slide, x - Inches(0.02), line_y - connector_h,
                     Inches(0.04), connector_h, fill_color=accent)
        else:
            add_rect(slide, x - Inches(0.02), line_y + dot_size/2,
                     Inches(0.04), connector_h, fill_color=accent)

        # 文字區塊
        box_w = step * 0.85
        box_x = x - box_w / 2
        if is_above:
            # 日期（小字，在標題上方）
            add_text(slide, node.get("date",""), box_x, line_y - connector_h - Inches(0.9),
                     box_w, Inches(0.3), font_size=11, color=accent,
                     align=PP_ALIGN.CENTER, font_name="Calibri")
            # 標題
            add_text(slide, node.get("title",""), box_x, line_y - connector_h - Inches(0.6),
                     box_w, Inches(0.45), font_size=13, bold=True, color=DARK,
                     align=PP_ALIGN.CENTER)
            # 說明
            add_text(slide, node.get("description",""), box_x, line_y - connector_h - Inches(0.18),
                     box_w, Inches(0.38), font_size=11, color=RGBColor(100,100,100),
                     align=PP_ALIGN.CENTER)
        else:
            add_text(slide, node.get("date",""), box_x,
                     line_y + connector_h + Inches(0.0),
                     box_w, Inches(0.3), font_size=11, color=accent,
                     align=PP_ALIGN.CENTER, font_name="Calibri")
            add_text(slide, node.get("title",""), box_x,
                     line_y + connector_h + Inches(0.3),
                     box_w, Inches(0.45), font_size=13, bold=True, color=DARK,
                     align=PP_ALIGN.CENTER)
            add_text(slide, node.get("description",""), box_x,
                     line_y + connector_h + Inches(0.72),
                     box_w, Inches(0.38), font_size=11, color=RGBColor(100,100,100),
                     align=PP_ALIGN.CENTER)


def _build_process_steps(slide, s: dict):
    """步驟流程圖（橫向，含箭頭）"""
    accent  = parse_color(s.get("accent_color", "#0091D5"))
    primary = parse_color(s.get("brand_color",  "#1B4F9B"))
    DARK    = RGBColor(64,64,64)
    WHITE   = RGBColor(255,255,255)
    title   = s.get("slide_title","")
    steps   = s.get("steps", [])   # list of {title, description}

    add_text(slide, title, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
             font_size=24, bold=True, color=DARK)

    n = min(len(steps), 5)
    if n == 0:
        return

    # 每個步驟的寬度
    total_w  = Inches(12.5)
    arrow_w  = Inches(0.4)
    step_w   = (total_w - arrow_w * (n - 1)) / n
    start_x  = Inches(0.4)
    box_top  = Inches(1.5)
    box_h    = Inches(4.5)

    for i, step in enumerate(steps[:n]):
        x = start_x + (step_w + arrow_w) * i

        # 步驟底色矩形
        fill = primary if i == n - 1 else RGBColor(0xE8,0xF4,0xFD)
        text_c = WHITE if i == n - 1 else DARK
        add_rect(slide, x, box_top, step_w, box_h, fill_color=fill)

        # 數字徽章（頂部）
        badge_size = Inches(0.55)
        add_rect(slide, x + step_w/2 - badge_size/2, box_top + Inches(0.2),
                 badge_size, badge_size, fill_color=accent)
        add_text(slide, str(i+1),
                 x + step_w/2 - badge_size/2, box_top + Inches(0.2),
                 badge_size, badge_size,
                 font_size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font_name="Calibri")

        # 步驟標題
        add_text(slide, step.get("title",""), x + Inches(0.1), box_top + Inches(0.9),
                 step_w - Inches(0.2), Inches(0.6),
                 font_size=14, bold=True, color=text_c, align=PP_ALIGN.CENTER)

        # 步驟說明
        add_text(slide, step.get("description",""), x + Inches(0.15), box_top + Inches(1.6),
                 step_w - Inches(0.3), Inches(2.5),
                 font_size=12, color=text_c, align=PP_ALIGN.CENTER)

        # 箭頭（非最後一個）
        if i < n - 1:
            ax = x + step_w + Inches(0.05)
            ay = box_top + box_h/2 - Inches(0.15)
            add_rect(slide, ax, ay, arrow_w - Inches(0.1), Inches(0.3),
                     fill_color=accent)


def _build_three_columns(slide, s: dict):
    """三欄支柱 / 三大亮點"""
    accent  = parse_color(s.get("accent_color", "#0091D5"))
    primary = parse_color(s.get("brand_color",  "#1B4F9B"))
    DARK    = RGBColor(64,64,64)
    WHITE   = RGBColor(255,255,255)
    title   = s.get("slide_title","")
    columns = s.get("columns", [])   # list of {icon_label, title, points:[str]}

    add_text(slide, title, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             font_size=24, bold=True, color=DARK)

    col_w   = Inches(3.8)
    col_gap = Inches(0.4)
    col_top = Inches(1.1)
    col_h   = Inches(5.8)
    starts  = [Inches(0.4), Inches(0.4) + col_w + col_gap,
               Inches(0.4) + (col_w + col_gap) * 2]

    for i, col in enumerate(columns[:3]):
        x      = starts[i]
        is_mid = (i == 1) and s.get("highlight_middle", False)
        bg     = primary if is_mid else RGBColor(0xF4,0xF7,0xFF)
        txt_c  = WHITE   if is_mid else DARK

        add_rect(slide, x, col_top, col_w, col_h, fill_color=bg)

        # Icon 圓形佔位
        ic_size = Inches(0.7)
        add_rect(slide, x + col_w/2 - ic_size/2, col_top + Inches(0.25),
                 ic_size, ic_size, fill_color=accent)
        add_text(slide, col.get("icon_label","★"),
                 x + col_w/2 - ic_size/2, col_top + Inches(0.25),
                 ic_size, ic_size,
                 font_size=20, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font_name="Calibri")

        # 欄位標題
        add_text(slide, col.get("title",""), x + Inches(0.1), col_top + Inches(1.1),
                 col_w - Inches(0.2), Inches(0.55),
                 font_size=16, bold=True, color=txt_c, align=PP_ALIGN.CENTER)

        # 重點條列
        for j, pt in enumerate(col.get("points",[])[:5]):
            add_rect(slide, x + Inches(0.2), col_top + Inches(1.85) + Inches(0.6)*j,
                     Inches(0.15), Inches(0.15), fill_color=accent)
            add_text(slide, pt, x + Inches(0.45), col_top + Inches(1.78) + Inches(0.6)*j,
                     col_w - Inches(0.55), Inches(0.5),
                     font_size=12, color=txt_c)


def _build_quote(slide, s: dict):
    """金句引言 / 客戶見證"""
    accent  = parse_color(s.get("accent_color",  "#0091D5"))
    primary = parse_color(s.get("brand_color",   "#1B4F9B"))
    is_dark = s.get("dark_background", False)
    bg      = primary if is_dark else RGBColor(255,255,255)
    txt_c   = RGBColor(255,255,255) if is_dark else RGBColor(40,40,40)
    DARK    = RGBColor(64,64,64)

    # 背景
    add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=bg)

    # 超大引號裝飾
    add_text(slide, "\u201C", Inches(0.5), Inches(0.3), Inches(3), Inches(2.0),
             font_size=120, bold=True, color=accent, font_name="Calibri")

    # 引言正文
    quote_text = s.get("quote_text", "")
    add_text(slide, quote_text, Inches(1.2), Inches(1.8),
             Inches(10.5), Inches(3.0),
             font_size=28, bold=True, color=txt_c, align=PP_ALIGN.CENTER)

    # 說話者資訊（右下）
    speaker = s.get("speaker_name","")
    role    = s.get("speaker_role","")
    if speaker:
        add_text(slide, f"— {speaker}", Inches(8.0), Inches(5.5),
                 Inches(4.8), Inches(0.5),
                 font_size=16, bold=True, color=accent, align=PP_ALIGN.RIGHT,
                 font_name="Calibri")
    if role:
        add_text(slide, role, Inches(8.0), Inches(6.0),
                 Inches(4.8), Inches(0.4),
                 font_size=13, color=txt_c, align=PP_ALIGN.RIGHT,
                 font_name="Calibri")


def _build_cta(slide, s: dict):
    """行動呼籲 / 聯絡資訊頁"""
    primary = parse_color(s.get("brand_color",  "#1B4F9B"))
    accent  = parse_color(s.get("accent_color", "#0091D5"))
    is_dark = s.get("dark_background", True)
    bg      = primary if is_dark else RGBColor(255,255,255)
    txt_c   = RGBColor(255,255,255) if is_dark else RGBColor(40,40,40)

    add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=bg)

    # 主標題
    add_text(slide, s.get("headline",""), Inches(0.5), Inches(0.4),
             Inches(12), Inches(0.8),
             font_size=30, bold=True, color=txt_c, align=PP_ALIGN.CENTER)

    # QR Code 佔位（左側）
    qr_size = Inches(3.2)
    qr_x    = Inches(1.5)
    qr_y    = Inches(1.8)
    add_rect(slide, qr_x, qr_y, qr_size, qr_size,
             fill_color=RGBColor(255,255,255))
    add_rect(slide, qr_x + Inches(0.15), qr_y + Inches(0.15),
             qr_size - Inches(0.3), qr_size - Inches(0.3),
             fill_color=RGBColor(30,30,30))
    add_text(slide, "QR", qr_x, qr_y, qr_size, qr_size,
             font_size=32, bold=True, color=RGBColor(255,255,255),
             align=PP_ALIGN.CENTER, font_name="Calibri")

    add_text(slide, s.get("qr_label","掃描取得資料"),
             qr_x, qr_y + qr_size + Inches(0.1), qr_size, Inches(0.4),
             font_size=12, color=txt_c, align=PP_ALIGN.CENTER)

    # 聯絡資訊（右側）
    info_x = Inches(5.8)
    info_y = Inches(2.0)
    for item in s.get("contact_items",[])[:4]:
        add_rect(slide, info_x, info_y + Inches(0.1),
                 Inches(0.25), Inches(0.25), fill_color=accent)
        add_text(slide, item, info_x + Inches(0.4), info_y,
                 Inches(7.0), Inches(0.45),
                 font_size=15, color=txt_c, font_name="Calibri")
        info_y += Inches(0.7)

    # 底部 Logo
    logo = s.get("logo_text","")
    if logo:
        add_text(slide, logo, Inches(0.5), Inches(6.9),
                 Inches(12), Inches(0.4),
                 font_size=10, bold=True, color=RGBColor(200,220,255),
                 align=PP_ALIGN.CENTER, font_name="Calibri")


def _build_speaker_profile(slide, s: dict):
    """講者 / 團隊成員介紹"""
    accent  = parse_color(s.get("accent_color", "#0091D5"))
    primary = parse_color(s.get("brand_color",  "#1B4F9B"))
    DARK    = RGBColor(40,40,40)
    WHITE   = RGBColor(255,255,255)

    # 左側照片佔位
    photo_w = Inches(4.5)
    add_rect(slide, Inches(0.4), Inches(0.8), photo_w, Inches(5.5),
             fill_color=RGBColor(180,190,200))
    add_text(slide, s.get("speaker_name",""), Inches(0.4), Inches(6.1),
             photo_w, Inches(0.6),
             font_size=22, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(slide, s.get("speaker_role",""), Inches(0.4), Inches(6.65),
             photo_w, Inches(0.4),
             font_size=14, color=accent, align=PP_ALIGN.CENTER, font_name="Calibri")

    # 右側內容
    rx = Inches(5.4)
    # 裝飾線
    add_rect(slide, rx, Inches(0.8), Inches(7.5), Inches(0.05), fill_color=accent)

    # 介紹標題
    add_text(slide, s.get("intro_heading","關於講者"), rx, Inches(1.0),
             Inches(7.0), Inches(0.5),
             font_size=13, bold=True, color=accent, font_name="Calibri")

    # 專長條列
    for i, item in enumerate(s.get("expertise",[])[:5]):
        add_rect(slide, rx, Inches(1.65) + Inches(0.75)*i,
                 Inches(0.18), Inches(0.35), fill_color=accent)
        add_text(slide, item, rx + Inches(0.3), Inches(1.6) + Inches(0.75)*i,
                 Inches(7.0), Inches(0.5),
                 font_size=14, color=DARK)

    # 底部裝飾線
    add_rect(slide, rx, Inches(7.1), Inches(7.5), Inches(0.05), fill_color=accent)


def _build_icon_grid(slide, s: dict):
    """圖標網格（2×2 或 2×3）"""
    accent  = parse_color(s.get("accent_color", "#0091D5"))
    primary = parse_color(s.get("brand_color",  "#1B4F9B"))
    DARK    = RGBColor(40,40,40)
    WHITE   = RGBColor(255,255,255)
    title   = s.get("slide_title","")
    cells   = s.get("cells", [])   # list of {icon_label, title, description}

    add_text(slide, title, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             font_size=24, bold=True, color=DARK)

    rows = 2
    cols = min(3, -(-len(cells[:6]) // rows))   # ceiling division
    cell_w = Inches(12.0) / cols
    cell_h = Inches(5.8) / rows
    start_x, start_y = Inches(0.6), Inches(1.1)

    for idx, cell in enumerate(cells[:6]):
        row = idx // cols
        col = idx %  cols
        cx  = start_x + cell_w * col
        cy  = start_y + cell_h * row

        # 格子邊框
        add_rect(slide, cx, cy, cell_w - Inches(0.1), cell_h - Inches(0.1),
                 fill_color=RGBColor(248,249,255),
                 border_color=RGBColor(220,225,240))

        # Icon 佔位
        ic = Inches(0.6)
        add_rect(slide, cx + cell_w/2 - ic/2, cy + Inches(0.2), ic, ic,
                 fill_color=accent)
        add_text(slide, cell.get("icon_label","●"),
                 cx + cell_w/2 - ic/2, cy + Inches(0.2), ic, ic,
                 font_size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font_name="Calibri")

        # 標題
        add_text(slide, cell.get("title",""), cx + Inches(0.1), cy + Inches(0.95),
                 cell_w - Inches(0.2), Inches(0.45),
                 font_size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)

        # 說明
        add_text(slide, cell.get("description",""), cx + Inches(0.15), cy + Inches(1.45),
                 cell_w - Inches(0.3), cell_h - Inches(1.65),
                 font_size=12, color=RGBColor(80,80,80), align=PP_ALIGN.CENTER)


def _build_big_title(slide, s: dict):
    """大標題頁：留白底 + 巨大主句 + 可選副文字"""
    bg_color    = parse_color(s.get("bg_color", "#FFFFFF"))
    text_color  = parse_color(s.get("text_color", "#1A1A1A"))
    accent_color = parse_color(s.get("accent_color", "#1B4F9B"))
    logo_text   = s.get("logo_text", "")

    add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=bg_color)

    # 頂部細色條
    add_rect(slide, 0, 0, SLIDE_WIDTH, Inches(0.08), fill_color=accent_color)

    # Logo
    if logo_text:
        add_text(slide, logo_text, Inches(0.30), Inches(0.15),
                 Inches(6), Inches(0.40), font_size=9, bold=True,
                 color=accent_color, font_name="Calibri")

    main_title = s.get("main_title", "")
    sub_text   = s.get("sub_text", "")

    # 主標題垂直置中
    title_top = Inches(2.50) if not sub_text else Inches(1.90)
    add_text(slide, main_title, Inches(0.80), title_top,
             Inches(11.70), Inches(2.50),
             font_size=int(s.get("title_size", 60)), bold=True,
             color=text_color, align=PP_ALIGN.CENTER)

    if sub_text:
        add_text(slide, sub_text, Inches(1.50), Inches(4.80),
                 Inches(10.30), Inches(1.20),
                 font_size=24, bold=False,
                 color=RGBColor(100, 100, 100), align=PP_ALIGN.CENTER)


def _build_full_image(slide, s: dict):
    """全圖頁：滿版圖片（或深色佔位）+ 少量文字疊加"""
    overlay_color = parse_color(s.get("overlay_color", "#0D1B2A"))
    text_color    = RGBColor(255, 255, 255)
    logo_text     = s.get("logo_text", "")
    image_path    = s.get("image_path")

    if image_path and Path(image_path).exists():
        slide.shapes.add_picture(image_path, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
        # 半透明遮罩 — 用深色矩形近似
        add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT,
                 fill_color=parse_color(s.get("overlay_color", "#0D1B2A")))
    else:
        add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=overlay_color)

    if logo_text:
        add_text(slide, logo_text, Inches(0.30), Inches(0.15),
                 Inches(6), Inches(0.40), font_size=9, bold=True,
                 color=RGBColor(200, 200, 200), font_name="Calibri")

    main_title = s.get("main_title", "")
    caption    = s.get("caption", "")

    add_text(slide, main_title, Inches(0.80), Inches(2.50),
             Inches(11.70), Inches(2.50),
             font_size=int(s.get("title_size", 54)), bold=True,
             color=text_color, align=PP_ALIGN.CENTER)

    if caption:
        add_text(slide, caption, Inches(0.80), Inches(5.40),
                 Inches(11.70), Inches(0.80),
                 font_size=18, color=RGBColor(200, 200, 200),
                 align=PP_ALIGN.CENTER)


def _build_data_chart(slide, s: dict):
    """數據圖表佔位頁：標題 + KPI 數字區塊 + 圖表說明"""
    accent_color = parse_color(s.get("accent_color", "#1B4F9B"))
    logo_text    = s.get("logo_text", "")

    add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=RGBColor(255, 255, 255))
    add_rect(slide, 0, 0, SLIDE_WIDTH, Inches(0.08), fill_color=accent_color)

    if logo_text:
        add_text(slide, logo_text, Inches(0.30), Inches(0.15),
                 Inches(6), Inches(0.40), font_size=9, bold=True,
                 color=accent_color, font_name="Calibri")

    slide_title = s.get("slide_title", "")
    add_text(slide, slide_title, Inches(0.50), Inches(0.35),
             Inches(12.30), Inches(0.60),
             font_size=24, bold=True, color=RGBColor(40, 40, 40))

    # KPI 區塊（最多 3 個）
    kpis = s.get("kpis", [])[:3]
    if kpis:
        kpi_w = Inches(3.80)
        for i, kpi in enumerate(kpis):
            lx = Inches(0.30 + i * 4.40)
            add_rect(slide, lx, Inches(1.10), kpi_w, Inches(2.00),
                     fill_color=parse_color(s.get("kpi_bg", "#EEF3FA")))
            add_text(slide, kpi.get("value", "—"), lx + Inches(0.20), Inches(1.20),
                     kpi_w - Inches(0.40), Inches(1.00),
                     font_size=42, bold=True, color=accent_color, align=PP_ALIGN.CENTER)
            add_text(slide, kpi.get("label", ""), lx + Inches(0.20), Inches(2.30),
                     kpi_w - Inches(0.40), Inches(0.60),
                     font_size=13, color=RGBColor(80, 80, 80), align=PP_ALIGN.CENTER)

    # 圖表佔位說明
    chart_note = s.get("chart_note", "[ 圖表佔位：請在 Google Slides / PowerPoint 插入圖表 ]")
    add_rect(slide, Inches(0.30), Inches(3.30), Inches(12.70), Inches(3.80),
             fill_color=RGBColor(245, 245, 245),
             border_color=RGBColor(200, 200, 200))
    add_text(slide, chart_note, Inches(0.30), Inches(4.80),
             Inches(12.70), Inches(1.00),
             font_size=14, color=RGBColor(160, 160, 160), align=PP_ALIGN.CENTER)

    source = s.get("source", "")
    if source:
        add_text(slide, f"資料來源：{source}", Inches(0.30), Inches(6.95),
                 Inches(12.70), Inches(0.40),
                 font_size=10, color=RGBColor(150, 150, 150))


def _build_two_columns(slide, s: dict):
    """雙欄對比頁：左欄 vs 右欄，支援 Before/After 或優缺點對比"""
    accent_color = parse_color(s.get("accent_color", "#1B4F9B"))
    logo_text    = s.get("logo_text", "")

    add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=RGBColor(255, 255, 255))
    add_rect(slide, 0, 0, SLIDE_WIDTH, Inches(0.08), fill_color=accent_color)

    if logo_text:
        add_text(slide, logo_text, Inches(0.30), Inches(0.15),
                 Inches(6), Inches(0.40), font_size=9, bold=True,
                 color=accent_color, font_name="Calibri")

    slide_title = s.get("slide_title", "")
    add_text(slide, slide_title, Inches(0.50), Inches(0.35),
             Inches(12.30), Inches(0.60),
             font_size=24, bold=True, color=RGBColor(40, 40, 40))

    left  = s.get("left_column",  {})
    right = s.get("right_column", {})

    col_w = Inches(6.00)
    gap   = Inches(0.33)

    for col, lx, header_color in [
        (left,  Inches(0.30), parse_color(s.get("left_color",  "#2D415F"))),
        (right, Inches(6.96), parse_color(s.get("right_color", "#C0392B"))),
    ]:
        # 欄標題背景
        add_rect(slide, lx, Inches(1.10), col_w, Inches(0.55), fill_color=header_color)
        add_text(slide, col.get("title", ""), lx + Inches(0.20), Inches(1.13),
                 col_w - Inches(0.40), Inches(0.50),
                 font_size=16, bold=True, color=RGBColor(255, 255, 255))

        # 欄內容框
        add_rect(slide, lx, Inches(1.70), col_w, Inches(5.50),
                 fill_color=parse_color(s.get("col_bg", "#F7F9FC")),
                 border_color=RGBColor(220, 220, 220))

        bullets = col.get("bullets", [])
        for j, b in enumerate(bullets[:6]):
            add_text(slide, f"• {b}", lx + Inches(0.25), Inches(1.85 + j * 0.75),
                     col_w - Inches(0.50), Inches(0.65),
                     font_size=14, color=RGBColor(50, 50, 50))


def _build_right_img_left_points(slide, s: dict):
    """右圖左重點：左側條列重點，右側圖片"""
    accent_color = parse_color(s.get("accent_color", "#1B4F9B"))
    logo_text    = s.get("logo_text", "")

    add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=RGBColor(255, 255, 255))
    add_rect(slide, 0, 0, SLIDE_WIDTH, Inches(0.08), fill_color=accent_color)

    if logo_text:
        add_text(slide, logo_text, Inches(0.30), Inches(0.15),
                 Inches(6), Inches(0.40), font_size=9, bold=True,
                 color=accent_color, font_name="Calibri")

    slide_title = s.get("slide_title", "")
    add_text(slide, slide_title, Inches(0.50), Inches(0.35),
             Inches(7.00), Inches(0.60),
             font_size=22, bold=True, color=RGBColor(40, 40, 40))

    # 右側圖片佔位
    img_path = s.get("image_path")
    img_l, img_t = Inches(6.80), Inches(1.00)
    img_w, img_h = Inches(6.30), Inches(6.20)
    if img_path and Path(img_path).exists():
        slide.shapes.add_picture(img_path, int(img_l), int(img_t), int(img_w), int(img_h))
    else:
        add_rect(slide, img_l, img_t, img_w, img_h,
                 fill_color=RGBColor(0x1A, 0x34, 0x5C))
        add_text(slide, s.get("image_caption", "[ 圖片 ]"),
                 img_l + Inches(0.20), img_t + Inches(2.80),
                 img_w - Inches(0.40), Inches(0.60),
                 font_size=14, color=RGBColor(200, 200, 200), align=PP_ALIGN.CENTER)

    # 左側重點
    bullets = s.get("bullets", [])
    for i, b in enumerate(bullets[:6]):
        add_rect(slide, Inches(0.30), Inches(1.20 + i * 0.90),
                 Inches(0.08), Inches(0.50), fill_color=accent_color)
        add_text(slide, b, Inches(0.55), Inches(1.18 + i * 0.90),
                 Inches(5.90), Inches(0.75),
                 font_size=15, color=RGBColor(50, 50, 50))


def _build_cycle_diagram(slide, s: dict):
    """環狀循環圖：中央標題 + 外圍 3-6 個階段（以矩形近似）"""
    accent_color = parse_color(s.get("accent_color", "#1B4F9B"))
    logo_text    = s.get("logo_text", "")

    add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=RGBColor(255, 255, 255))
    add_rect(slide, 0, 0, SLIDE_WIDTH, Inches(0.08), fill_color=accent_color)

    if logo_text:
        add_text(slide, logo_text, Inches(0.30), Inches(0.15),
                 Inches(6), Inches(0.40), font_size=9, bold=True,
                 color=accent_color, font_name="Calibri")

    slide_title = s.get("slide_title", "")
    add_text(slide, slide_title, Inches(0.50), Inches(0.35),
             Inches(12.30), Inches(0.60),
             font_size=22, bold=True, color=RGBColor(40, 40, 40))

    # 中央核心橢圓（用矩形近似）
    center_label = s.get("center_label", "")
    cx, cy = Inches(5.50), Inches(3.15)
    cw, ch = Inches(2.33), Inches(1.40)
    add_rect(slide, cx, cy, cw, ch, fill_color=accent_color)
    add_text(slide, center_label, cx + Inches(0.10), cy + Inches(0.20),
             cw - Inches(0.20), ch - Inches(0.40),
             font_size=16, bold=True, color=RGBColor(255, 255, 255),
             align=PP_ALIGN.CENTER)

    # 外圍階段（最多 6 個）
    import math
    stages = s.get("stages", [])[:6]
    n = len(stages)
    if n == 0:
        return
    radius_x, radius_y = Inches(4.20), Inches(2.60)
    center_x = Inches(6.665)
    center_y = Inches(3.85)
    stage_w, stage_h = Inches(2.20), Inches(0.90)

    stage_colors = [
        parse_color(c) for c in s.get("stage_colors", [
            "#2D415F","#1B4F9B","#0091D5","#00B4D8","#48CAE4","#90E0EF"
        ])
    ]

    for i, stage in enumerate(stages):
        angle = math.pi / 2 + (2 * math.pi * i / n)
        sx = center_x + radius_x * math.cos(angle) - stage_w / 2
        sy = center_y - radius_y * math.sin(angle) - stage_h / 2
        c  = stage_colors[i % len(stage_colors)]
        add_rect(slide, sx, sy, stage_w, stage_h, fill_color=c)
        add_text(slide, stage.get("title",""), sx + Inches(0.10), sy + Inches(0.10),
                 stage_w - Inches(0.20), stage_h - Inches(0.20),
                 font_size=13, bold=True, color=RGBColor(255,255,255),
                 align=PP_ALIGN.CENTER)


def _build_pyramid_funnel(slide, s: dict):
    """金字塔/漏斗圖：3-5 層遞減矩形條（由下而上寬→窄）"""
    accent_color = parse_color(s.get("accent_color", "#1B4F9B"))
    logo_text    = s.get("logo_text", "")
    mode         = s.get("mode", "pyramid")  # "pyramid" or "funnel"

    add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=RGBColor(255, 255, 255))
    add_rect(slide, 0, 0, SLIDE_WIDTH, Inches(0.08), fill_color=accent_color)

    if logo_text:
        add_text(slide, logo_text, Inches(0.30), Inches(0.15),
                 Inches(6), Inches(0.40), font_size=9, bold=True,
                 color=accent_color, font_name="Calibri")

    slide_title = s.get("slide_title", "")
    add_text(slide, slide_title, Inches(0.50), Inches(0.35),
             Inches(12.30), Inches(0.60),
             font_size=22, bold=True, color=RGBColor(40, 40, 40))

    tiers = s.get("tiers", [])[:5]
    n = len(tiers)
    if n == 0:
        return

    default_colors = ["#1B4F9B","#2E6FC0","#4A90D9","#7BB3E8","#B3D4F5"]
    max_w    = Inches(8.00)
    min_w    = Inches(2.00)
    bar_h    = Inches(0.88)
    gap      = Inches(0.08)
    total_h  = n * bar_h + (n - 1) * gap
    start_y  = (SLIDE_HEIGHT - total_h) / 2 + Inches(0.40)
    center_x = SLIDE_WIDTH / 2

    for i, tier in enumerate(tiers):
        if mode == "funnel":
            w = max_w - (max_w - min_w) * i / max(n - 1, 1)
        else:  # pyramid (wide at bottom)
            w = min_w + (max_w - min_w) * i / max(n - 1, 1)
        lx = center_x - w / 2
        ty = start_y + i * (bar_h + gap)
        c  = parse_color(tier.get("color", default_colors[i % len(default_colors)]))
        add_rect(slide, lx, ty, w, bar_h, fill_color=c)
        add_text(slide, tier.get("title",""), lx + Inches(0.15), ty + Inches(0.15),
                 w - Inches(0.30), bar_h - Inches(0.30),
                 font_size=15, bold=True, color=RGBColor(255,255,255),
                 align=PP_ALIGN.CENTER)
        desc = tier.get("desc","")
        if desc:
            add_text(slide, desc, center_x + max_w/2 + Inches(0.20), ty + Inches(0.15),
                     Inches(3.00), bar_h - Inches(0.30),
                     font_size=12, color=RGBColor(80, 80, 80))


def _build_image_gallery(slide, s: dict):
    """多圖畫廊：2-4 張圖片拼貼 + 圖說"""
    accent_color = parse_color(s.get("accent_color", "#1B4F9B"))
    logo_text    = s.get("logo_text", "")

    add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=RGBColor(255, 255, 255))
    add_rect(slide, 0, 0, SLIDE_WIDTH, Inches(0.08), fill_color=accent_color)

    if logo_text:
        add_text(slide, logo_text, Inches(0.30), Inches(0.15),
                 Inches(6), Inches(0.40), font_size=9, bold=True,
                 color=accent_color, font_name="Calibri")

    slide_title = s.get("slide_title", "")
    add_text(slide, slide_title, Inches(0.50), Inches(0.35),
             Inches(12.30), Inches(0.60),
             font_size=22, bold=True, color=RGBColor(40, 40, 40))

    photos = s.get("photos", [])[:4]
    n      = len(photos)
    if n == 0:
        return

    cols    = min(n, 2) if n <= 2 else min(n, 4)
    rows    = (n + cols - 1) // cols
    gap     = Inches(0.15)
    total_w = SLIDE_WIDTH - Inches(0.60)
    total_h = Inches(5.80)
    cell_w  = (total_w - gap * (cols - 1)) / cols
    cell_h  = (total_h - gap * (rows - 1)) / rows

    placeholder_colors = [
        RGBColor(0x1A, 0x34, 0x5C),
        RGBColor(0x2E, 0x6F, 0xC0),
        RGBColor(0x0A, 0x4B, 0x7C),
        RGBColor(0x12, 0x56, 0x8A),
    ]

    for i, photo in enumerate(photos):
        row = i // cols
        col = i % cols
        lx  = Inches(0.30) + col * (cell_w + gap)
        ty  = Inches(1.10) + row * (cell_h + gap)

        img_path = photo.get("image_path")
        if img_path and Path(img_path).exists():
            slide.shapes.add_picture(img_path, int(lx), int(ty), int(cell_w), int(cell_h))
        else:
            add_rect(slide, lx, ty, cell_w, cell_h,
                     fill_color=placeholder_colors[i % len(placeholder_colors)])

        caption = photo.get("caption","")
        if caption:
            add_text(slide, caption, lx, ty + cell_h - Inches(0.45),
                     cell_w, Inches(0.40),
                     font_size=11, color=RGBColor(255,255,255),
                     align=PP_ALIGN.CENTER)


def _build_video_demo(slide, s: dict):
    """影片展示頁：筆電 Mockup 框 + 播放鍵 + 引導標題"""
    accent_color = parse_color(s.get("accent_color", "#1B4F9B"))
    logo_text    = s.get("logo_text", "")

    add_rect(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=RGBColor(18, 18, 18))

    if logo_text:
        add_text(slide, logo_text, Inches(0.30), Inches(0.15),
                 Inches(6), Inches(0.40), font_size=9, bold=True,
                 color=RGBColor(180, 180, 180), font_name="Calibri")

    slide_title = s.get("slide_title", "")
    add_text(slide, slide_title, Inches(0.50), Inches(0.30),
             Inches(12.30), Inches(0.65),
             font_size=22, bold=True, color=RGBColor(255, 255, 255),
             align=PP_ALIGN.CENTER)

    # 筆電外框（深灰矩形）
    frame_l, frame_t = Inches(1.80), Inches(1.10)
    frame_w, frame_h = Inches(9.70), Inches(5.80)
    add_rect(slide, frame_l, frame_t, frame_w, frame_h,
             fill_color=RGBColor(40, 40, 40),
             border_color=RGBColor(80, 80, 80))

    # 螢幕佔位
    screen_l = frame_l + Inches(0.30)
    screen_t = frame_t + Inches(0.30)
    screen_w = frame_w - Inches(0.60)
    screen_h = frame_h - Inches(1.00)
    img_path  = s.get("image_path")
    if img_path and Path(img_path).exists():
        slide.shapes.add_picture(img_path, int(screen_l), int(screen_t),
                                  int(screen_w), int(screen_h))
    else:
        add_rect(slide, screen_l, screen_t, screen_w, screen_h,
                 fill_color=RGBColor(0, 0, 0))

    # 播放鍵（三角近似——用細矩形組合）
    play_cx = screen_l + screen_w / 2
    play_cy = screen_t + screen_h / 2
    add_rect(slide, play_cx - Inches(0.35), play_cy - Inches(0.35),
             Inches(0.70), Inches(0.70),
             fill_color=accent_color)
    add_text(slide, "▶", play_cx - Inches(0.35), play_cy - Inches(0.35),
             Inches(0.70), Inches(0.70),
             font_size=22, bold=True,
             color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)

    # 底部說明文字
    caption = s.get("caption", "")
    if caption:
        add_text(slide, caption, Inches(0.50), Inches(6.80),
                 Inches(12.30), Inches(0.55),
                 font_size=13, color=RGBColor(160, 160, 160),
                 align=PP_ALIGN.CENTER)


# ── 版型路由表 ────────────────────────────────────────────────────────────────

LAYOUT_BUILDERS = {
    # 基礎版型
    "cover":             _build_cover_or_closing,
    "closing":           _build_cover_or_closing,
    "contents":          _build_contents,
    "section_divider":   _build_section_divider,
    "bullet_points":     _build_bullet_points,
    "table":             _build_table,
    # 進階版型 - 流程與結構
    "timeline":          _build_timeline,
    "process_steps":     _build_process_steps,
    "cycle_diagram":     _build_cycle_diagram,
    "pyramid_funnel":    _build_pyramid_funnel,
    # 進階版型 - 資訊圖表與視覺
    "big_title":         _build_big_title,
    "full_image":        _build_full_image,
    "data_chart":        _build_data_chart,
    "two_columns":       _build_two_columns,
    "three_columns":     _build_three_columns,
    "icon_grid":         _build_icon_grid,
    "image_gallery":     _build_image_gallery,
    "right_img_left_points": _build_right_img_left_points,
    # 進階版型 - 信任感與互動
    "speaker_profile":   _build_speaker_profile,
    "quote":             _build_quote,
    "video_demo":        _build_video_demo,
    "cta":               _build_cta,
}

# ── 公開入口 ──────────────────────────────────────────────────────────────────

def build_presentation_from_plan(plan: list[dict], output_path) -> str:
    """
    根據 slide plan 生成 PPTX。

    Args:
        plan: 投影片清單，每項為 dict，必須含 "layout" 欄位
        output_path: 輸出路徑（str 或 Path）

    Returns:
        絕對路徑字串
    """
    prs = Presentation()
    prs.slide_width  = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank = prs.slide_layouts[6]

    for s in plan:
        layout  = s.get("layout","").lower().replace("-","_")
        builder = LAYOUT_BUILDERS.get(layout)
        slide   = prs.slides.add_slide(blank)

        if builder:
            builder(slide, s)
        else:
            add_text(slide,
                     f"[未支援的版型: {layout}]\n支援：{', '.join(LAYOUT_BUILDERS)}",
                     Inches(1), Inches(3), Inches(11), Inches(1.5),
                     font_size=18, color=RGBColor(200,0,0))

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return str(out.resolve())
